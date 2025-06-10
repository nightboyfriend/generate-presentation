from pptx import Presentation
import os

def check_template(template_path="templates/template.pptx"):
    if not os.path.exists(template_path):
        print(f"Ошибка: Файл шаблона {template_path} не найден")
        return

    try:
        prs = Presentation(template_path)
        print(f"Шаблон загружен: {template_path}")
        print(f"Количество слайдов: {len(prs.slides)}")

        if len(prs.slides) < 2:
            print("Ошибка: Шаблон должен содержать минимум 2 слайда (титульный и основной)")
            return

        title_slide = prs.slides[0]
        if title_slide.shapes.title:
            print("Титульный слайд: Найден плейсхолдер заголовка")
        else:
            print("Предупреждение: На титульном слайде отсутствует плейсхолдер заголовка")

        if len(prs.slides) > 1:
            content_slide = prs.slides[1]
            placeholders = [shape for shape in content_slide.shapes if shape.is_placeholder]
            title_placeholder = any(p.placeholder_format.type == 1 for p in placeholders)  
            body_placeholder = any(p.placeholder_format.type == 7 for p in placeholders)  
            picture_placeholder = any(p.placeholder_format.type == 18 for p in placeholders)
            print("Основной слайд:")
            print(f"  Плейсхолдер заголовка: {'найден' if title_placeholder else 'не найден'}")
            print(f"  Плейсхолдер текста: {'найден' if body_placeholder else 'не найден'}")
            print(f"  Плейсхолдер изображения: {'найден' if picture_placeholder else 'не найден'}")
        else:
            print("Ошибка: Основной слайд отсутствует")

        if len(prs.slides) > 2:
            print("Последний слайд: Найден")
        else:
            print("Предупреждение: Последний слайд отсутствует, он не будет добавлен")

    except Exception as e:
        print(f"Ошибка при проверке шаблона: {e}")

if __name__ == "__main__":
    check_template()