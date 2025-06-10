from pptx import Presentation
from pptx.util import Pt, Inches
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE_TYPE
import os
from typing import List, Dict
from loguru import logger

def create_presentation(slide_count: int) -> Presentation:
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(5.625)
    return prs

def load_template(template_path: str = "templates/template.pptx") -> Presentation:
    """Загружает шаблон презентации."""
    if not os.path.exists(template_path):
        raise FileNotFoundError(f"Шаблон {template_path} не найден")
    return Presentation(template_path)

def add_template_title_slide(prs: Presentation, topic: str) -> None:
    """Добавляет титульный слайд из шаблона, заменяя текст темы."""
    if len(prs.slides) < 1:
        raise ValueError("Шаблон не содержит титульный слайд")
    slide = prs.slides[0]  # Предполагаем, что первый слайд — титульный
    title = slide.shapes.title
    if title:
        title.text = topic
        logger.debug(f"Титульный слайд: Текст заменён на '{topic}'")
    else:
        logger.warning("На титульном слайде отсутствует плейсхолдер заголовка")

def add_template_content_slide(prs: Presentation, data: Dict, slide_index: int) -> None:
    """Добавляет основной слайд из шаблона, заменяя заголовок, описание и изображение."""
    if len(prs.slide_layouts) < 2:
        raise ValueError("Шаблон не содержит макет для основного слайда")
    slide_layout = prs.slide_layouts[1]  # Предполагаем, что второй макет — для основного слайда
    slide = prs.slides.add_slide(slide_layout)
    logger.debug(f"Добавлен основной слайд {slide_index} с макетом {slide_layout.name}")

    # Замена заголовка
    title = slide.shapes.title
    if title and 'zagolovok' in data and data['zagolovok']:
        title.text = data['zagolovok']
        logger.debug(f"Слайд {slide_index}: Заголовок заменён на '{data['zagolovok']}'")
    else:
        logger.warning(f"Слайд {slide_index}: Плейсхолдер заголовка не найден или заголовок отсутствует")

    # Замена описания
    content_placeholder = None
    for shape in slide.placeholders:
        if shape.placeholder_format.type == 7:  # Body/Content placeholder
            content_placeholder = shape
            break
    if content_placeholder and 'opisanie' in data and data['opisanie']:
        content_placeholder.text_frame.clear()
        p = content_placeholder.text_frame.add_paragraph()
        p.text = data['opisanie']
        logger.debug(f"Слайд {slide_index}: Описание заменено на '{data['opisanie']}'")
    else:
        logger.warning(f"Слайд {slide_index}: Плейсхолдер текста не найден или описание отсутствует")

    # Поиск плейсхолдера изображения
    picture_placeholder = None
    for shape in slide.placeholders:
        if shape.placeholder_format.type == 18:  # Picture placeholder
            picture_placeholder = shape
            break

    # Добавление изображения
    if 'photo' in data and data['photo'] and os.path.exists(data['photo']):
        try:
            if picture_placeholder:
                # Используем плейсхолдер изображения
                picture_placeholder.insert_picture(data['photo'])
                logger.debug(f"Слайд {slide_index}: Изображение вставлено в плейсхолдер")
            else:
                # Если плейсхолдера нет, добавляем изображение по координатам
                slide.shapes.add_picture(data['photo'], Inches(6), Inches(1), width=Inches(3))
                logger.debug(f"Слайд {slide_index}: Изображение добавлено по координатам Inches(6), Inches(1)")
        except Exception as e:
            logger.error(f"Слайд {slide_index}: Ошибка добавления изображения '{data.get('photo')}': {e}")
    else:
        logger.warning(f"Слайд {slide_index}: Изображение не найдено или путь некорректен: {data.get('photo', 'Отсутствует')}")

def add_template_final_slide(prs: Presentation, template_prs: Presentation) -> None:
    """Копирует последний слайд из шаблона, если он существует."""
    if len(template_prs.slides) >= 3:
        slide_layout = template_prs.slide_layouts[2]  # Предполагаем, что третий слайд — последний
        prs.slides.add_slide(slide_layout)
        logger.debug("Последний слайд добавлен из шаблона")
    else:
        logger.warning("Последний слайд в шаблоне отсутствует, не добавлен")

def add_slide(prs: Presentation, data: Dict) -> None:
    """Добавляет слайд для обычного режима."""
    slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(slide_layout)

    if 'zagolovok' in data and data['zagolovok']:
        title = slide.shapes.title
        if title:
            title.text = data['zagolovok']
            title.text_frame.paragraphs[0].font.name = 'Arial'
            title.text_frame.paragraphs[0].font.size = Pt(28)
            title.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
            title.top = Inches(0.5)
            title.width = Inches(10)
            title.text_frame.paragraphs[0].font.bold = True

    if 'opisanie' in data and data['opisanie']:
        if len(slide.placeholders) > 1:
            content = slide.placeholders[1]
            content.text_frame.clear()
            p = content.text_frame.add_paragraph()
            p.text = data['opisanie']
            p.font.name = 'Arial'
            p.font.size = Pt(14)
            p.alignment = PP_ALIGN.LEFT
            content.left = Inches(0.5)
            content.top = Inches(1)
            content.width = Inches(5)
            content.height = Inches(4.5)

    if 'photo' in data and data['photo'] and os.path.exists(data['photo']):
        try:
            slide.shapes.add_picture(data['photo'], Inches(6), Inches(1), width=Inches(3))
        except Exception as e:
            print(f"Ошибка добавления фото для слайда '{data.get('zagolovok', 'Неизвестно')}': {e}")

def generate_presentation(data: List[Dict], slide_count: int, output_path: str, topic: str, template_mode: bool = False) -> None:
    if not data and not topic:
        print("Нет данных или темы для генерации презентации")
        return
    
    if template_mode:
        prs = load_template()
        # Добавляем титульный слайд
        add_template_title_slide(prs, topic)
        # Добавляем основные слайды (дублируем второй слайд из шаблона)
        for i, slide_data in enumerate(data[:slide_count - 2]):  # -2 для титульного и последнего
            add_template_content_slide(prs, slide_data, i + 1)
        # Добавляем последний слайд
        add_template_final_slide(prs, prs)
    else:
        prs = create_presentation(slide_count)
        # Добавляем титульный слайд
        if topic:
            slide_layout = prs.slide_layouts[0]
            slide = prs.slides.add_slide(slide_layout)
            title = slide.shapes.title
            if title:
                title.text = topic
                title.text_frame.paragraphs[0].font.name = 'Arial'
                title.text_frame.paragraphs[0].font.size = Pt(36)
                title.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
                title.text_frame.paragraphs[0].font.bold = True
                title.left = Inches(0)
                title.top = Inches(2)
                title.width = Inches(10)
                title.height = Inches(1.5)
        # Добавляем остальные слайды
        for i, slide_data in enumerate(data[:slide_count - 1]):
            add_slide(prs, slide_data)
    
    try:
        prs.save(output_path)
        print(f"Презентация сохранена в {output_path}")
    except Exception as e:
        print(f"Ошибка сохранения презентации: {e}")